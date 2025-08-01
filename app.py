import os
import glob
import shutil
import uuid
from flask import Flask, request, render_template, send_from_directory, url_for
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import numpy as np
from skimage import io
from skimage.color import rgb2gray
from skimage.filters import sobel

# --- FLASK APP CONFIGURATION ---
app = Flask(__name__)
# Create a dedicated temp folder for uploaded files to keep them separate
app.config['UPLOADS_FOLDER'] = os.path.join('temp_uploads') 
app.config['CROPPED_FOLDER_BASE'] = os.path.join('cropped')
app.config['OUTPUT_FOLDER'] = os.path.join('output')
app.config['ALLOWED_EXTENSIONS'] = {'png', 'jpg', 'jpeg'}

# --- Ensure all necessary folders exist ---
os.makedirs(app.config['UPLOADS_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
os.makedirs(app.config['CROPPED_FOLDER_BASE'], exist_ok=True)

def allowed_file(filename):
    """Checks if the uploaded file has an allowed extension."""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# --- IMAGE PROCESSING AND PPT LOGIC (Adapted from your script) ---

def autocrop_image(input_path, output_path, threshold=0.01):
    """Analyzes an image, finds the content boundaries, crops it, and saves the result."""
    try:
        image = io.imread(input_path)
        if image.shape[2] == 4:
            image = image[:, :, :3]
            
        grayscale = rgb2gray(image)
        edges = sobel(grayscale) > threshold
        
        coords = np.argwhere(edges)
        if coords.size == 0:
            shutil.copy(input_path, output_path)
            return

        y0, x0 = coords.min(axis=0)
        y1, x1 = coords.max(axis=0)
        
        padding = 20
        y0 = max(0, y0 - padding)
        x0 = max(0, x0 - padding)
        y1 = min(image.shape[0], y1 + padding)
        x1 = min(image.shape[1], x1 + padding)

        cropped_image = image[y0:y1, x0:x1]
        io.imsave(output_path, cropped_image)
    except Exception as e:
        print(f"Could not crop {os.path.basename(input_path)}: {e}. Copying original.")
        shutil.copy(input_path, output_path)

def create_presentation(question_files, background_image_path, output_pptx_path):
    """Main function to generate the PowerPoint presentation."""
    
    # 1. CROP IMAGES
    request_id = str(uuid.uuid4())
    cropped_folder_path = os.path.join(app.config['CROPPED_FOLDER_BASE'], request_id)
    os.makedirs(cropped_folder_path, exist_ok=True)
    
    cropped_image_paths = []
    print("--- Starting Auto-Cropping ---")
    for i, original_path in enumerate(question_files):
        question_num = i + 1
        cropped_file_path = os.path.join(cropped_folder_path, f'cropped_q{question_num}.png')
        print(f"Processing '{os.path.basename(original_path)}'")
        autocrop_image(original_path, cropped_file_path)
        cropped_image_paths.append(cropped_file_path)
    print("--- Auto-Cropping Complete ---\n")

    # 2. CREATE THE PRESENTATION
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]

    def add_standard_slide(title_text):
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        return slide

    def add_question_slide(question_num, image_path):
        slide = add_standard_slide(f"Question {question_num}")
        max_width, max_height = Inches(5.0), Inches(7.0)
        try:
            im = Image.open(image_path)
            img_width_px, img_height_px = im.size
            aspect_ratio = img_width_px / img_height_px
            if (max_height * aspect_ratio) > max_width:
                new_width, new_height = max_width, max_width / aspect_ratio
            else:
                new_width, new_height = max_height * aspect_ratio, max_height
            left = Inches(0.5) + (max_width - new_width) / 2
            top = Inches(1.5) + (max_height - new_height) / 2
            slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
        except Exception as e:
            print(f"Could not add image {os.path.basename(image_path)} to slide: {e}")

    print("--- Starting Presentation Creation ---")
    add_standard_slide("Advanced Problems in Functions and Calculus\n\nA Practice Set")
    add_standard_slide("Welcome & Objectives")
    
    for i, cropped_img_path in enumerate(cropped_image_paths):
        question_num = i + 1
        print(f"Adding slides for Question {question_num}...")
        add_question_slide(question_num, cropped_img_path)
        add_standard_slide(f"Solution for Question {question_num}")

    add_standard_slide("Thank You & Q/A")

    # 3. SAVE THE FILE AND CLEAN UP
    prs.save(output_pptx_path)
    shutil.rmtree(cropped_folder_path) # Clean up the temporary cropped images folder
    print(f"\nSuccess! Presentation saved as '{os.path.basename(output_pptx_path)}'")


# --- FLASK ROUTES ---

@app.route('/')
def index():
    """Renders the main upload page."""
    return render_template('index.html')

@app.route('/create', methods=['POST'])
def create_ppt_route():
    """Handles file uploads and initiates the PPT creation process."""
    if 'questions' not in request.files or 'background' not in request.files:
        return "Error: Both question images and a background image are required.", 400

    question_uploads = request.files.getlist('questions')
    background_upload = request.files['background']

    if not question_uploads or not background_upload or background_upload.filename == '':
        return "Error: No files selected.", 400

    # --- START OF FIX ---
    # Save uploaded files temporarily with UNIQUE names
    
    question_paths = []
    temp_files_to_delete = [] # Keep a list of all temp files to delete later

    for file in question_uploads:
        if file and allowed_file(file.filename):
            # Generate a unique filename to prevent collisions
            original_filename = secure_filename(file.filename)
            extension = original_filename.rsplit('.', 1)[1]
            unique_filename = f"{uuid.uuid4().hex}.{extension}"
            
            path = os.path.join(app.config['UPLOADS_FOLDER'], unique_filename)
            file.save(path)
            question_paths.append(path)
            temp_files_to_delete.append(path) # Add to our cleanup list

    # Do the same for the background image
    bg_original_filename = secure_filename(background_upload.filename)
    bg_extension = bg_original_filename.rsplit('.', 1)[1]
    bg_unique_filename = f"{uuid.uuid4().hex}.{bg_extension}"
    
    bg_path = os.path.join(app.config['UPLOADS_FOLDER'], bg_unique_filename)
    background_upload.save(bg_path)
    temp_files_to_delete.append(bg_path) # Add background to cleanup list
    
    # --- END OF FIX ---

    output_filename = f"presentation_{uuid.uuid4().hex}.pptx"
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)

    # Run the presentation creation logic
    create_presentation(question_paths, bg_path, output_path)
    
    # Clean up all the unique temporary files from the uploads folder
    for path in temp_files_to_delete:
        try:
            os.remove(path)
        except OSError as e:
            print(f"Error deleting file {path}: {e}")

    return render_template('result.html', filename=output_filename)

@app.route('/download/<filename>')
def download_file(filename):
    """Serves the generated file for download."""
    return send_from_directory(app.config['OUTPUT_FOLDER'], filename, as_attachment=True)


if __name__ == '__main__':
    app.run(debug=True)