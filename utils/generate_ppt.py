import os
import glob
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import numpy as np
from skimage import io
from skimage.color import rgb2gray
from skimage.filters import sobel

# --- CONFIGURATION ---
# --- You can change these folder names if you like ---
ORIGINAL_IMAGES_FOLDER = 'original_questions'
CROPPED_IMAGES_FOLDER = 'cropped_questions'
BACKGROUND_IMAGE = 'background.jpg'
OUTPUT_FILENAME = 'Final_Automated_Presentation.pptx'

# --- HELPER FUNCTION: AUTO-CROP IMAGE ---
def autocrop_image(input_path, output_path, threshold=0.01):
    """Analyzes an image, finds the content boundaries, crops it, and saves the result."""
    try:
        image = io.imread(input_path)
        if image.shape[2] == 4:  # Handle transparency by taking only RGB channels
            image = image[:, :, :3]
            
        grayscale = rgb2gray(image)
        edges = sobel(grayscale) > threshold
        
        coords = np.argwhere(edges)
        if coords.size == 0:  # If no content is found, copy the original
            print(f"  -> No content detected in {os.path.basename(input_path)}. Copying original.")
            shutil.copy(input_path, output_path)
            return

        # Find the bounding box of the content
        y0, x0 = coords.min(axis=0)
        y1, x1 = coords.max(axis=0)
        
        # Add a small padding for aesthetics
        padding = 15
        y0 = max(0, y0 - padding)
        x0 = max(0, x0 - padding)
        y1 = min(image.shape[0], y1 + padding)
        x1 = min(image.shape[1], x1 + padding)

        cropped_image = image[y0:y1, x0:x1]
        io.imsave(output_path, cropped_image)
    except Exception as e:
        print(f"  -> Could not crop {os.path.basename(input_path)} due to error: {e}. Copying original instead.")
        shutil.copy(input_path, output_path)

# --- MAIN SCRIPT LOGIC ---

# 1. VALIDATE FOLDER AND FILE STRUCTURE
if not os.path.exists(ORIGINAL_IMAGES_FOLDER):
    print(f"Error: The required folder '{ORIGINAL_IMAGES_FOLDER}' was not found. Please create it and add your question images.")
    exit()

if not os.path.exists(BACKGROUND_IMAGE):
    print(f"Error: The background image '{BACKGROUND_IMAGE}' was not found in the main directory.")
    exit()

original_file_paths = glob.glob(os.path.join(ORIGINAL_IMAGES_FOLDER, '*.png'))
if not original_file_paths:
    print(f"Error: No .png files found in the '{ORIGINAL_IMAGES_FOLDER}' folder.")
    exit()

print(f"Found {len(original_file_paths)} question images to process.")

# 2. PERFORM AUTO-CROPPING
if os.path.exists(CROPPED_IMAGES_FOLDER):
    shutil.rmtree(CROPPED_IMAGES_FOLDER) # Start fresh every time
os.makedirs(CROPPED_IMAGES_FOLDER)
print(f"Created a fresh '{CROPPED_IMAGES_FOLDER}' folder.")

print("\n--- Starting Auto-Cropping Process ---")
for i, original_path in enumerate(original_file_paths):
    question_num = i + 1
    cropped_file = os.path.join(CROPPED_IMAGES_FOLDER, f'cropped_q{question_num}.png')
    print(f"Processing '{os.path.basename(original_path)}' -> 'cropped_q{question_num}.png'")
    autocrop_image(original_path, cropped_file)
print("--- Auto-Cropping Complete ---\n")

# 3. CREATE THE PRESENTATION
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)
blank_slide_layout = prs.slide_layouts[6]

def add_standard_slide(title_text):
    slide = prs.slides.add_slide(blank_slide_layout)
    pic = slide.shapes.add_picture(BACKGROUND_IMAGE, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    return slide

def add_question_slide(question_num, image_path):
    slide = add_standard_slide(f"Question {question_num}")
    max_width, max_height = Inches(5.0), Inches(7.0)
    try:
        im = Image.open(image_path)
        img_width_px, img_height_px = im.size
        aspect_ratio = img_width_px / img_height_px
        if (max_height * aspect_ratio) > max_width:
            new_width, new_height = max_width, max_width / aspect_ratio
        else:
            new_width, new_height = max_height * aspect_ratio, max_height
        left = Inches(0.5) + (max_width - new_width) / 2
        top = Inches(1.5) + (max_height - new_height) / 2
        slide.shapes.add_picture(image_path, left, top, width=new_width, height=new_height)
    except Exception as e:
        print(f"  -> Could not add image {os.path.basename(image_path)} to slide: {e}")

print("--- Starting Presentation Creation ---")

# Add Introductory Slides
add_standard_slide("Advanced Problems in Functions and Calculus\n\nA Practice Set")
add_standard_slide("Welcome & Objectives")
add_standard_slide("How to Use This Deck")
add_standard_slide("Table of Contents")

# Add Question and Solution Slides
num_questions_found = len(original_file_paths)
for i in range(num_questions_found):
    question_num = i + 1
    cropped_img_path = os.path.join(CROPPED_IMAGES_FOLDER, f'cropped_q{question_num}.png')
    print(f"Adding slides for Question {question_num}...")
    add_question_slide(question_num, cropped_img_path)
    add_standard_slide(f"Solution for Question {question_num} (Part 1)")
    add_standard_slide(f"Solution for Question {question_num} (Part 2)")

# Add Concluding Slides
add_standard_slide("End of Problem Set")
add_standard_slide("Key Takeaways & Formulas")
add_standard_slide("Thank You & Q/A")

# 4. SAVE THE FILE
prs.save(OUTPUT_FILENAME)
print(f"\nSuccess! Presentation complete. File saved as '{OUTPUT_FILENAME}'")`